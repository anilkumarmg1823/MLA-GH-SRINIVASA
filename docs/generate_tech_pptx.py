"""Generate purchase cost PowerPoint: Now / Monthly / 6 months."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLUE = RGBColor(0x00, 0x1D, 0x56)
BLUE_DEEP = RGBColor(0x00, 0x14, 0x38)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x33, 0x41, 0x55)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
TOTAL = 6


def solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, rgb):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    solid(s, rgb)
    return s


def set_run(run, text, size=18, bold=False, color=WHITE):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_text(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)
    return box


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.15), BLUE_DEEP)
    add_rect(slide, 0, Inches(1.15), W, Inches(0.08), GOLD)
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
             [(title, 28, True, WHITE)])
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.35),
                 [(subtitle, 14, False, GOLD)])


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(10), Inches(0.3),
             [("Kudligi MLA Portal — Setup now + Monthly + 6 months", 10, False, MUTED)])
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             [(f"{page} / {TOTAL}", 10, False, MUTED)], align=PP_ALIGN.RIGHT)


def bullet_block(slide, l, t, w, h, items, size=16, color=SLATE):
    lines = [(f"•  {item}", size, False, color) for item in items]
    add_text(slide, l, t, w, h, lines)


# ===== 1 Title =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, W, H, BLUE_DEEP)
add_rect(s, 0, Inches(5.9), W, Inches(1.6), BLUE)
add_rect(s, 0, Inches(5.9), W, Inches(0.1), GOLD)
add_text(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(2.5),
         [("Payment Plan — What to Buy", 36, True, WHITE),
          ("Pay now (setup / deploy)  →  Then monthly  →  6-month total", 20, False, GOLD)])
add_text(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.9),
         [("Kudligi MLA Digital Portal", 16, True, WHITE),
          ("Infrastructure purchase estimates (INR) — for client discussion", 13, False, GOLD)])

# ===== 2 AWS what to buy =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, W, H, WHITE)
header_bar(s, "AWS — Tier to purchase", "Small DB + file storage for uploads")
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.0))
solid(sh, LIGHT)
add_rect(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.55), BLUE)
add_text(s, Inches(0.6), Inches(1.58), Inches(5.8), Inches(0.4),
         [("RDS Database", 18, True, GOLD)])
bullet_block(s, Inches(0.65), Inches(2.25), Inches(5.7), Inches(4), [
    "Tier: db.t3.micro",
    "Storage: 20 GB (start)",
    "Later if needed: t3.small + 50 GB",
    "Monthly after free: ≈ ₹1,300 – ₹2,200",
    "With free tier / credits: often ~₹0 first months",
], size=15)
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.5), Inches(6.0), Inches(5.0))
solid(sh, LIGHT)
add_rect(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(0.55), GOLD)
add_text(s, Inches(7.05), Inches(1.58), Inches(5.6), Inches(0.4),
         [("S3 Files (uploads)", 18, True, BLUE_DEEP)])
bullet_block(s, Inches(7.1), Inches(2.25), Inches(5.5), Inches(4), [
    "Plan: 50 GB → 100 GB Year 1",
    "50 GB ≈ ₹40 – ₹100 / month",
    "100 GB ≈ ₹200 / month",
    "API host (Railway/VPS): ≈ ₹400 – ₹1,300 / month",
    "AWS total typical: ≈ ₹2,000 – ₹4,500 / month after free",
], size=15)
footer(s, 2)

# ===== 3 PAY NOW =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, W, H, WHITE)
header_bar(s, "Pay NOW — Setup / go-live purchase", "One-time amount to start deployment")
add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
         [("Client can give this amount now for instant deploy setup.", 15, False, SLATE)])

rows = [
    ("Domain (.in / .com) — 1 year", "₹700 – ₹1,500"),
    ("Fast2SMS wallet top-up (OTP start)", "₹1,000 – ₹2,000"),
    ("DLT / OTP template (if needed)", "₹0 – ₹1,000"),
    ("First AWS / hosting month (or free-tier ≈ ₹0)", "₹0 – ₹3,000"),
    ("Deploy / install setup (go-live work)", "₹5,000 – ₹15,000*"),
]
for i, (a, b) in enumerate(rows):
    t = Inches(1.9 + i * 0.75)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), t, Inches(12.3), Inches(0.65))
    solid(sh, LIGHT)
    add_text(s, Inches(0.8), t + Inches(0.15), Inches(8), Inches(0.4), [(a, 15, True, BLUE)])
    add_text(s, Inches(9.2), t + Inches(0.15), Inches(3.3), Inches(0.4), [(b, 15, True, BLUE_DEEP)])

add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.8),
         [("*Deploy/install = one-time setup to put website + API live (not monthly). Adjust as per your quote.", 12, False, MUTED),
          ("PAY NOW TOTAL (typical):  ≈ ₹7,000 – ₹22,000", 18, True, BLUE)],)
footer(s, 3)

# ===== 4 MONTHLY after =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, W, H, WHITE)
header_bar(s, "Then MONTHLY — after setup", "Recurring charges every month (lower than setup)")

# Two columns: early months vs normal
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.5), Inches(6.2), Inches(4.8))
solid(sh, LIGHT)
add_rect(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.55), GOLD)
add_text(s, Inches(0.6), Inches(1.58), Inches(5.8), Inches(0.4),
         [("Months 1–3 (free tier / soft)", 16, True, BLUE_DEEP)])
bullet_block(s, Inches(0.65), Inches(2.3), Inches(5.7), Inches(3.5), [
    "AWS: ≈ ₹0 – ₹2,000 (credits / free)",
    "Fast2SMS OTP: ≈ ₹500 – ₹1,200",
    "AI posters (optional): ₹0 – ₹500",
    "",
    "MONTHLY ≈ ₹500 – ₹3,500",
], size=15)

sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.5), Inches(6.0), Inches(4.8))
solid(sh, LIGHT)
add_rect(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(0.55), BLUE)
add_text(s, Inches(7.05), Inches(1.58), Inches(5.6), Inches(0.4),
         [("Normal production month", 16, True, GOLD)])
bullet_block(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(3.5), [
    "AWS RDS + S3 + API: ≈ ₹2,500 – ₹6,000",
    "Fast2SMS OTP: ≈ ₹500 – ₹1,200",
    "AI (optional): ≈ ₹100 – ₹500",
    "",
    "MONTHLY ≈ ₹3,000 – ₹8,000",
], size=15)
footer(s, 4)

# ===== 5 SMS + AI + Domain reminder (compact purchase lines) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, W, H, WHITE)
header_bar(s, "Other purchases (quick)", "SMS · Domain · Optional AI")
rows = [
    ("Fast2SMS", "₹0.25 / OTP (+GST)", "100 users ≈ ₹550 / month"),
    ("Domain", "₹700 – ₹1,500 / year", "Buy once at setup"),
    ("AI posters (optional)", "Grok ≈ ₹2 / poster", "50 posters ≈ ₹85 – ₹500 / month"),
]
for i, (a, b, c) in enumerate(rows):
    t = Inches(1.8 + i * 1.4)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), t, Inches(12.1), Inches(1.2))
    solid(sh, LIGHT)
    add_rect(s, Inches(0.6), t, Inches(0.12), Inches(1.2), GOLD)
    add_text(s, Inches(1.0), t + Inches(0.25), Inches(3.5), Inches(0.7), [(a, 18, True, BLUE)])
    add_text(s, Inches(4.8), t + Inches(0.25), Inches(3.8), Inches(0.7), [(b, 15, False, SLATE)])
    add_text(s, Inches(8.8), t + Inches(0.25), Inches(3.5), Inches(0.7), [(c, 15, True, BLUE_DEEP)])
footer(s, 5)

# ===== 6 TOTAL: Now + Monthly + 6 months =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, 0, 0, W, H, WHITE)
header_bar(s, "TOTAL — Now + Monthly + 6 Months", "Simple numbers to tell the client")

# Three cards
cards = [
    (0.4, GOLD, BLUE_DEEP, "PAY NOW", "Setup / deploy", "≈ ₹7,000 – ₹22,000", "Domain + SMS wallet + first month + deploy setup"),
    (4.7, BLUE, GOLD, "EVERY MONTH", "After go-live", "≈ ₹3,000 – ₹8,000", "AWS + OTP (+ AI if used). Soft months can be less."),
    (9.0, BLUE_DEEP, GOLD, "6 MONTHS", "Now + 6 months run", "≈ ₹25,000 – ₹70,000", "Pay now + ~6 × monthly (see note below)"),
]
for l, accent, title_c, h1, h2, amount, note in cards:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(1.45), Inches(4.0), Inches(3.6))
    solid(sh, LIGHT)
    add_rect(s, Inches(l), Inches(1.45), Inches(4.0), Inches(0.7), accent)
    add_text(s, Inches(l + 0.2), Inches(1.55), Inches(3.6), Inches(0.5),
             [(h1, 16, True, title_c)])
    add_text(s, Inches(l + 0.2), Inches(2.35), Inches(3.6), Inches(0.4),
             [(h2, 13, False, MUTED)])
    add_text(s, Inches(l + 0.2), Inches(2.9), Inches(3.6), Inches(0.8),
             [(amount, 20, True, BLUE)])
    add_text(s, Inches(l + 0.2), Inches(3.8), Inches(3.6), Inches(1.0),
             [(note, 12, False, SLATE)])

add_rect(s, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.35), BLUE_DEEP)
add_text(s, Inches(0.7), Inches(5.45), Inches(12), Inches(1.1),
         [("Example mid-range: Pay now ₹12,000 + (₹5,000 × 6) = ≈ ₹42,000 for 6 months", 16, True, GOLD),
          ("Soft launch first 2–3 months can be cheaper if AWS free tier / credits are still active.", 13, False, WHITE)])
footer(s, 6)

out = os.path.join(os.path.dirname(__file__), "Kudligi-MLA-Tech-Deployment.pptx")
prs.save(out)
print("Saved:", out)
